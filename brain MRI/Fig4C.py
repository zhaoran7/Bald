#!/usr/bin/env python3
from pathlib import Path
import re, shutil, tempfile, urllib.request
import numpy as np
import pandas as pd
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.colors import LinearSegmentedColormap
import nibabel as nib
from nilearn import plotting

from enigmatoolbox.utils.parcellation import parcel_to_surface
from enigmatoolbox.plotting import plot_cortical, plot_subcortical
from enigmatoolbox.datasets import load_fsa5


def patch_enigma_vtk_cells():
    import numpy as _np
    from vtk import vtkCellArray
    from vtk.util.numpy_support import numpy_to_vtkIdTypeArray
    from enigmatoolbox.vtk_interface.wrappers.data_object import BSPolyData

    def _numpy2cells(cells):
        cells = _np.asarray(cells)
        if cells.ndim == 1:
            offset = 0; n_cells = 0
            while offset < cells.size:
                offset += int(cells[offset]) + 1; n_cells += 1
            vtk_cells = cells
        else:
            n_cells, n_points_cell = cells.shape
            vtk_cells = _np.empty((n_cells, n_points_cell + 1), dtype=_np.int64)
            vtk_cells[:, 0] = n_points_cell
            vtk_cells[:, 1:] = cells.astype(_np.int64, copy=False)
            vtk_cells = vtk_cells.ravel()
        ca = vtkCellArray()
        ca.SetCells(int(n_cells), numpy_to_vtkIdTypeArray(_np.ascontiguousarray(vtk_cells, dtype=_np.int64), deep=True))
        return ca

    BSPolyData._numpy2cells = staticmethod(_numpy2cells)


def patch_enigma_vtk_lookup_table():
    import numpy as _np
    import vtk as _vtk
    from vtk.util.numpy_support import numpy_to_vtk
    from enigmatoolbox.vtk_interface.wrappers.lookup_table import BSLookupTable

    def SetTable(self, table):
        if not hasattr(table, "GetClassName"):
            arr = _np.asarray(table)
            if arr.ndim == 1:
                arr = arr.reshape((-1, 4))
            if arr.dtype.kind == "f":
                finite = arr[_np.isfinite(arr)]
                if finite.size and float(finite.max()) <= 1.0:
                    arr = arr * 255.0
            arr = _np.clip(arr, 0, 255).astype(_np.uint8, copy=False)
            if arr.shape[1] == 3:
                arr = _np.c_[arr, _np.full((arr.shape[0], 1), 255, dtype=_np.uint8)]
            table = numpy_to_vtk(_np.ascontiguousarray(arr), deep=True, array_type=_vtk.VTK_UNSIGNED_CHAR)
            table.SetNumberOfComponents(4)
        self.VTKObject.SetTable(table)

    BSLookupTable.SetTable = SetTable


patch_enigma_vtk_cells()
patch_enigma_vtk_lookup_table()

root = Path("D:/bald/img") if Path("D:/bald/img").exists() else Path("/mnt/d/bald/img")
out = root / "plot"; out.mkdir(parents=True, exist_ok=True)
tmp = Path(tempfile.mkdtemp(prefix="Fig3C_"))
cell_dir = tmp / "_cells"; cell_dir.mkdir(exist_ok=True)

trait_order = ["Continuous", "M shape", "O shape", "U shape"]
row_order = [
    ("Cortical surface area", "cortex", "area"),
    ("Cortical thickness", "cortex", "thickness"),
    ("Cortical volume", "cortex", "volume"),
    ("Subcortical volume", "subcortex", "volume"),
    ("White matter FA", "white_matter", "FA"),
    ("White matter MD", "white_matter", "MD"),
]

# Layout controls.
# Final compact layout: narrow x-axis, compact y-axis, separated 4-view ENIGMA panels,
# and unsplit white-matter panels. Syntax-checked output file.
BRAIN_SCALE = 0.90
VIEW_N = 4
VIEW_GAP = 40
CELL_W = 700
CELL_H = 200
GAP_X = 8
GAP_Y = 2
LEFT_W = 200
RIGHT_W = 90
TOP_H = 50
BOTTOM_H = 8
BRAIN_BOX_W = 590
BRAIN_BOX_H = 120
BRAIN_TOP_PAD = 6
LABEL_Y_OFFSET = 23

PNG_DPI = 300
PPT_W_IN = 15.75
FONT_PT = 8
ZLIM = 12
COLOR_RANGE = (-ZLIM, ZLIM)
CELL_SIZE = (2400, 900)
CELL_SCALE = (4, 4)
WM_DPI = 700

plt.rcParams["font.family"] = "Arial"
plt.rcParams["pdf.fonttype"] = 42
plt.rcParams["ps.fonttype"] = 42

CMAP_NAME = "zscore_blue_gray_red"
try:
    cmap = LinearSegmentedColormap.from_list(CMAP_NAME, ["#2166AC", "#EDEDED", "#B2182B"], N=256)
    matplotlib.colormaps.register(cmap, name=CMAP_NAME, force=True)
except Exception:
    pass

DK = [
    "bankssts", "caudalanteriorcingulate", "caudalmiddlefrontal", "cuneus", "entorhinal",
    "fusiform", "inferiorparietal", "inferiortemporal", "isthmuscingulate", "lateraloccipital",
    "lateralorbitofrontal", "lingual", "medialorbitofrontal", "middletemporal", "parahippocampal",
    "paracentral", "parsopercularis", "parsorbitalis", "parstriangularis", "pericalcarine",
    "postcentral", "posteriorcingulate", "precentral", "precuneus", "rostralanteriorcingulate",
    "rostralmiddlefrontal", "superiorfrontal", "superiorparietal", "superiortemporal", "supramarginal",
    "frontalpole", "temporalpole", "transversetemporal", "insula"
]
CTX_ORDER = [("left", r) for r in DK] + [("right", r) for r in DK]
SCTX_ORDER = [
    ("left", "accumbens"), ("left", "amygdala"), ("left", "caudate"), ("left", "hippocampus"),
    ("left", "pallidum"), ("left", "putamen"), ("left", "thalamus"),
    ("right", "accumbens"), ("right", "amygdala"), ("right", "caudate"), ("right", "hippocampus"),
    ("right", "pallidum"), ("right", "putamen"), ("right", "thalamus"),
]


def clean(x):
    return re.sub(r"[^a-z0-9]", "", str(x).lower())


def read_phe():
    x = pd.read_csv(root / "res" / "mri_mr.tsv", sep="\t")
    x = x[(x["direction"] == "img_to_bald") & (x["status"] == "ok") & (x["fdr_mri272"] < 0.05)].copy()
    x["plot_score"] = (x["beta"] / x["se"]).clip(-ZLIM, ZLIM)
    x["region_clean"] = x["region"].map(clean)
    x["hemi"] = x["hemi"].astype(str).str.lower()
    return x


def label_text(x):
    x = re.sub(r"^(Cortical|Subcortical|White matter) (surface area|area|thickness|volume|FA|MD) ", "", str(x))
    x = re.sub(r"^(left|right|Left|Right) ", "", x)
    key = clean(x)
    m = {
        "bankssts":"Banks STS","caudalanteriorcingulate":"Caudal Anterior Cingulate","caudalmiddlefrontal":"Caudal Middle Frontal",
        "cuneus":"Cuneus","entorhinal":"Entorhinal","fusiform":"Fusiform","inferiorparietal":"Inferior Parietal","inferiortemporal":"Inferior Temporal",
        "isthmuscingulate":"Isthmus Cingulate","lateraloccipital":"Lateral Occipital","lateralorbitofrontal":"Lateral Orbitofrontal","lingual":"Lingual",
        "medialorbitofrontal":"Medial Orbitofrontal","middletemporal":"Middle Temporal","parahippocampal":"Parahippocampal","paracentral":"Paracentral",
        "parsopercularis":"Pars Opercularis","parsorbitalis":"Pars Orbitalis","parstriangularis":"Pars Triangularis","pericalcarine":"Pericalcarine",
        "postcentral":"Postcentral","posteriorcingulate":"Posterior Cingulate","precentral":"Precentral","precuneus":"Precuneus",
        "rostralanteriorcingulate":"Rostral Anterior Cingulate","rostralmiddlefrontal":"Rostral Middle Frontal","superiorfrontal":"Superior Frontal",
        "superiorparietal":"Superior Parietal","superiortemporal":"Superior Temporal","supramarginal":"Supramarginal","frontalpole":"Frontal Pole",
        "temporalpole":"Temporal Pole","transversetemporal":"Transverse Temporal","insula":"Insula","accumbensarea":"Accumbens Area","accumbens":"Accumbens Area",
        "amygdala":"Amygdala","caudate":"Caudate","hippocampus":"Hippocampus","pallidum":"Pallidum","putamen":"Putamen","thalamusproper":"Thalamus","thalamus":"Thalamus",
        "middlecerebellarpeduncle":"Middle Cerebellar Peduncle","pontinecrossingtract":"Pontine Crossing Tract","genuofcorpuscallosum":"Genu Of Corpus Callosum",
        "bodyofcorpuscallosum":"Body Of Corpus Callosum","spleniumofcorpuscallosum":"Splenium Of Corpus Callosum","fornix":"Fornix","corticospinaltract":"Corticospinal Tract",
        "mediallemniscus":"Medial Lemniscus","inferiorcerebellarpeduncle":"Inferior Cerebellar Peduncle","superiorcerebellarpeduncle":"Superior Cerebellar Peduncle",
        "cerebralpeduncle":"Cerebral Peduncle","anteriorlimbofinternalcapsule":"Anterior Limb Of Internal Capsule","posteriorlimbofinternalcapsule":"Posterior Limb Of Internal Capsule",
        "retrolenticularpartofinternalcapsule":"Retrolenticular Part Of Internal Capsule","anteriorcoronaradiata":"Anterior Corona Radiata","superiorcoronaradiata":"Superior Corona Radiata",
        "posteriorcoronaradiata":"Posterior Corona Radiata","posteriorthalamicradiation":"Posterior Thalamic Radiation","sagittalstratum":"Sagittal Stratum","externalcapsule":"External Capsule",
        "cingulumcingulategyrus":"Cingulum Cingulate Gyrus","cingulumhippocampus":"Cingulum Hippocampus","fornixcresandstriaterminalis":"Fornix Cres And Stria Terminalis",
        "superiorlongitudinalfasciculus":"Superior Longitudinal Fasciculus","superiorfrontooccipitalfasciculus":"Superior Fronto-Occipital Fasciculus",
        "inferiorfrontooccipitalfasciculus":"Inferior Fronto-Occipital Fasciculus","uncinatefasciculus":"Uncinate Fasciculus","tapetum":"Tapetum"
    }
    return m.get(key, " ".join(w.capitalize() for w in re.sub(r"([a-z])([A-Z])", r"\1 \2", x).split()))


def best_labels(matched, n=2):
    labs = []
    if matched is None or not len(matched):
        return labs
    for v in matched.sort_values("fdr_mri272")["standard_name"].tolist():
        lab = label_text(v)
        if lab not in labs:
            labs.append(lab)
        if len(labs) == n:
            break
    return labs


def empty_cell(file):
    Image.new("RGB", CELL_SIZE, "white").save(file, dpi=(PNG_DPI, PNG_DPI))
    return file


def cortical_array(d):
    vals = np.zeros(68, dtype=float); matched = []
    for i, (hemi, reg) in enumerate(CTX_ORDER):
        z = d[(d["hemi"] == hemi) & (d["region_clean"] == clean(reg))]
        if len(z):
            r = z.sort_values("fdr_mri272").iloc[0]
            vals[i] = float(r["plot_score"]); matched.append(r)
    return vals, pd.DataFrame(matched)


def subcortical_array(d):
    vals = np.zeros(14, dtype=float); matched = []
    for i, (hemi, reg) in enumerate(SCTX_ORDER):
        z = d[(d["hemi"] == hemi) & (d["region_clean"].str.contains(clean(reg), na=False))]
        if len(z):
            r = z.sort_values("fdr_mri272").iloc[0]
            vals[i] = float(r["plot_score"]); matched.append(r)
    return vals, pd.DataFrame(matched)


def make_vertex_array(vals):
    surf = np.asarray(parcel_to_surface(np.asarray(vals, dtype=float).reshape(-1), "aparc_fsa5", fill=0), dtype=float).squeeze().reshape(-1)
    lh, rh = load_fsa5()
    if surf.size != int(lh.n_points + rh.n_points):
        raise RuntimeError(f"aparc_fsa5 mapped array has length {surf.size}, but fsa5 surface needs {int(lh.n_points + rh.n_points)}.")
    return surf


def draw_cortex(d, file):
    vals, matched = cortical_array(d)
    if not np.any(vals != 0):
        return empty_cell(file), matched
    plot_cortical(array_name=make_vertex_array(vals), surface_name="fsa5", color_bar=False, cmap=CMAP_NAME,
                  color_range=COLOR_RANGE, nan_color=(0.93,0.93,0.93,1), background=(1,1,1),
                  size=CELL_SIZE, interactive=False, screenshot=True, filename=str(file),
                  transparent_bg=False, scale=CELL_SCALE, zoom=1.35)
    return file, matched


def draw_subcortex(d, file):
    vals, matched = subcortical_array(d)
    if not np.any(vals != 0):
        return empty_cell(file), matched
    plot_subcortical(array_name=vals, ventricles=False, color_bar=False, cmap=CMAP_NAME,
                     color_range=COLOR_RANGE, nan_color=(0.93,0.93,0.93,1), background=(1,1,1),
                     size=CELL_SIZE, interactive=False, screenshot=True, filename=str(file),
                     transparent_bg=False, scale=CELL_SCALE, zoom=1.25)
    return file, matched


def fetch_jhu():
    atlas_dir = root / "res" / "atlas"; atlas_dir.mkdir(parents=True, exist_ok=True)
    f = atlas_dir / "JHU-ICBM-labels-2mm.nii.gz"
    if not f.exists() or f.stat().st_size < 5000:
        urllib.request.urlretrieve("https://ftp.nmr.mgh.harvard.edu/pub/dist/freesurfer/tutorial_packages_centos6/centos6/freesurfer-fsl-matlab-Linux-centos6_x86_64-dev/freesurfer/fsl_507/data/atlases/JHU/JHU-ICBM-labels-2mm.nii.gz", f)
    labels = [
        "Background", "Middle cerebellar peduncle", "Pontine crossing tract", "Genu of corpus callosum", "Body of corpus callosum", "Splenium of corpus callosum", "Fornix",
        "Corticospinal tract R", "Corticospinal tract L", "Medial lemniscus R", "Medial lemniscus L", "Inferior cerebellar peduncle R", "Inferior cerebellar peduncle L",
        "Superior cerebellar peduncle R", "Superior cerebellar peduncle L", "Cerebral peduncle R", "Cerebral peduncle L", "Anterior limb of internal capsule R",
        "Anterior limb of internal capsule L", "Posterior limb of internal capsule R", "Posterior limb of internal capsule L", "Retrolenticular part of internal capsule R",
        "Retrolenticular part of internal capsule L", "Anterior corona radiata R", "Anterior corona radiata L", "Superior corona radiata R", "Superior corona radiata L",
        "Posterior corona radiata R", "Posterior corona radiata L", "Posterior thalamic radiation R", "Posterior thalamic radiation L", "Sagittal stratum R", "Sagittal stratum L",
        "External capsule R", "External capsule L", "Cingulum cingulate gyrus R", "Cingulum cingulate gyrus L", "Cingulum hippocampus R", "Cingulum hippocampus L",
        "Fornix cres and Stria terminalis R", "Fornix cres and Stria terminalis L", "Superior longitudinal fasciculus R", "Superior longitudinal fasciculus L",
        "Superior fronto-occipital fasciculus R", "Superior fronto-occipital fasciculus L", "Inferior fronto-occipital fasciculus R", "Inferior fronto-occipital fasciculus L",
        "Uncinate fasciculus R", "Uncinate fasciculus L"
    ]
    return type("Atlas", (), {"maps": nib.load(str(f)), "labels": labels})


def wm_label(row):
    x = str(row["region"]).lower(); side = "R" if row["hemi"] == "right" else "L"
    both = {"middle cerebellar peduncle":"Middle cerebellar peduncle", "pontine crossing tract":"Pontine crossing tract", "genu of corpus callosum":"Genu of corpus callosum", "body of corpus callosum":"Body of corpus callosum", "splenium of corpus callosum":"Splenium of corpus callosum", "fornix":"Fornix"}
    pair = {"corticospinal tract":"Corticospinal tract", "medial lemniscus":"Medial lemniscus", "inferior cerebellar peduncle":"Inferior cerebellar peduncle", "superior cerebellar peduncle":"Superior cerebellar peduncle", "cerebral peduncle":"Cerebral peduncle", "anterior limb of internal capsule":"Anterior limb of internal capsule", "posterior limb of internal capsule":"Posterior limb of internal capsule", "retrolenticular part of internal capsule":"Retrolenticular part of internal capsule", "anterior corona radiata":"Anterior corona radiata", "superior corona radiata":"Superior corona radiata", "posterior corona radiata":"Posterior corona radiata", "posterior thalamic radiation":"Posterior thalamic radiation", "sagittal stratum":"Sagittal stratum", "external capsule":"External capsule", "cingulum cingulate gyrus":"Cingulum cingulate gyrus", "cingulum hippocampus":"Cingulum hippocampus", "fornix cres and stria terminalis":"Fornix cres and Stria terminalis", "superior longitudinal fasciculus":"Superior longitudinal fasciculus", "superior fronto occipital fasciculus":"Superior fronto-occipital fasciculus", "inferior fronto occipital fasciculus":"Inferior fronto-occipital fasciculus", "uncinate fasciculus":"Uncinate fasciculus", "tapetum":"Tapetum"}
    for k, v in both.items():
        if k in x: return v
    for k, v in pair.items():
        if k in x: return f"{v} {side}"
    return None


def atlas_stat(d, atlas, label_fun):
    img = atlas.maps if hasattr(atlas.maps, "get_fdata") else nib.load(atlas.maps)
    dat = img.get_fdata(); outdat = np.zeros(dat.shape, dtype=float)
    labels = {str(v): i for i, v in enumerate(atlas.labels)}; matched = []
    for _, r in d.sort_values("fdr_mri272").iterrows():
        lab = label_fun(r)
        if lab not in labels: continue
        idx = labels[lab]
        if not np.any(outdat[dat == idx]):
            outdat[dat == idx] = r["plot_score"]
            rr = r.copy(); rr["atlas_label"] = lab; matched.append(rr)
    return nib.Nifti1Image(outdat, img.affine, img.header), pd.DataFrame(matched)


def draw_white_matter(d, atlas, file):
    stat, matched = atlas_stat(d, atlas, wm_label)
    if not np.any(stat.get_fdata() != 0):
        return empty_cell(file), matched
    disp = plotting.plot_glass_brain(stat, display_mode="lyrz", cmap=CMAP_NAME, symmetric_cbar=True,
                                     threshold=0.001, vmax=ZLIM, vmin=-ZLIM, colorbar=False,
                                     plot_abs=False, black_bg=False, title="")
    disp.savefig(str(file), dpi=WM_DPI); disp.close()
    return file, matched


def crop_white(im, pad=10):
    im = im.convert("RGB")
    bg = Image.new("RGB", im.size, "white")
    diff = ImageChops.difference(im, bg).convert("L")
    bbox = diff.point(lambda p: 255 if p > 8 else 0).getbbox()
    if bbox is None:
        return im
    l, t, r, b = bbox
    l = max(0, l - pad); t = max(0, t - pad); r = min(im.width, r + pad); b = min(im.height, b + pad)
    return im.crop((l, t, r, b))


def split_views(im, n=VIEW_N):
    im = crop_white(im, pad=6)
    w, h = im.size
    pieces = []
    for i in range(n):
        l = int(round(i * w / n))
        r = int(round((i + 1) * w / n))
        p = crop_white(im.crop((l, 0, r, h)), pad=3)
        pieces.append(p)
    return pieces


def make_view_strip(img_path):
    pieces = split_views(Image.open(img_path), VIEW_N)
    slot_w = max(1, (BRAIN_BOX_W - (VIEW_N - 1) * VIEW_GAP) // VIEW_N)
    strip = Image.new("RGB", (BRAIN_BOX_W, BRAIN_BOX_H), "white")
    for i, p in enumerate(pieces):
        s = min(slot_w / p.width, BRAIN_BOX_H / p.height)
        nw, nh = max(1, int(p.width * s)), max(1, int(p.height * s))
        p = p.resize((nw, nh), Image.Resampling.LANCZOS)
        slot_x = i * (slot_w + VIEW_GAP)
        x0 = slot_x + (slot_w - nw) // 2
        y0 = (BRAIN_BOX_H - nh) // 2
        strip.paste(p, (x0, y0))
    return strip


def make_single_panel(img_path):
    p = crop_white(Image.open(img_path), pad=6)
    s = min(BRAIN_BOX_W / p.width, BRAIN_BOX_H / p.height)
    nw, nh = max(1, int(p.width * s)), max(1, int(p.height * s))
    p = p.resize((nw, nh), Image.Resampling.LANCZOS)
    panel = Image.new("RGB", (BRAIN_BOX_W, BRAIN_BOX_H), "white")
    panel.paste(p, ((BRAIN_BOX_W - nw) // 2, (BRAIN_BOX_H - nh) // 2))
    return panel


def place_cell(canvas, img_path, x, y, split_four=True):
    cell = make_view_strip(img_path) if split_four else make_single_panel(img_path)
    px = x + (CELL_W - BRAIN_BOX_W) // 2
    py = y + BRAIN_TOP_PAD
    canvas.paste(cell, (px, py))


def stitch_no_text(cell_files):
    rows = [r for r in row_order if any((r[0], tr) in cell_files for tr in trait_order)]
    width = LEFT_W + 4 * CELL_W + 3 * GAP_X + RIGHT_W
    height = TOP_H + len(rows) * CELL_H + (len(rows) - 1) * GAP_Y + BOTTOM_H
    canvas = Image.new("RGB", (width, height), "white")
    for i, (row_lab, row_cls, _) in enumerate(rows):
        y = TOP_H + i * (CELL_H + GAP_Y)
        for j, tr in enumerate(trait_order):
            x = LEFT_W + j * (CELL_W + GAP_X)
            place_cell(canvas, cell_files[(row_lab, tr)], x, y, split_four=(row_cls != "white_matter"))
    cmap = plt.get_cmap(CMAP_NAME)
    from PIL import ImageDraw
    dr = ImageDraw.Draw(canvas)
    bar_h, bar_w = 230, 18
    x0, y0 = width - RIGHT_W + 38, 50
    for k in range(bar_h):
        dr.line((x0, y0 + k, x0 + bar_w, y0 + k), fill=tuple(int(255*c) for c in cmap(1-k/(bar_h-1))[:3]))
    dr.rectangle((x0, y0, x0 + bar_w, y0 + bar_h), outline="black", width=1)
    png = tmp / "Fig3C.png"
    canvas.save(png, dpi=(PNG_DPI, PNG_DPI))
    return png, (width, height), (x0, y0, bar_w, bar_h)


def add_textbox(slide, x, y, w, h, text, font_pt=8, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = "Arial"; r.font.size = Pt(font_pt)
    return tb


def px_to_ppt(val, total_px, total_in):
    return Inches(total_in * val / total_px)


def make_ppt(img_png, cell_labels, canvas_info, colorbar_info):
    canvas_w, canvas_h = canvas_info; x0, y0, bar_w, bar_h = colorbar_info
    slide_w = canvas_w / PNG_DPI if PPT_W_IN is None else PPT_W_IN
    slide_h = slide_w * canvas_h / canvas_w
    ppt = Presentation(); ppt.slide_width = Inches(slide_w); ppt.slide_height = Inches(slide_h)
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])
    slide.shapes.add_picture(str(img_png), 0, 0, width=Inches(slide_w), height=Inches(slide_h))

    for j, tr in enumerate(trait_order):
        x = LEFT_W + j * (CELL_W + GAP_X)
        add_textbox(slide, px_to_ppt(x, canvas_w, slide_w), px_to_ppt(12, canvas_h, slide_h),
                    px_to_ppt(CELL_W, canvas_w, slide_w), px_to_ppt(24, canvas_h, slide_h), tr, FONT_PT)

    rows = [r for r in row_order if any((r[0], tr) in cell_labels for tr in trait_order)]
    for i, (row_lab, _, _) in enumerate(rows):
        y = TOP_H + i * (CELL_H + GAP_Y)
        add_textbox(slide, px_to_ppt(18, canvas_w, slide_w), px_to_ppt(y + CELL_H/2 - 13, canvas_h, slide_h),
                    px_to_ppt(LEFT_W-36, canvas_w, slide_w), px_to_ppt(26, canvas_h, slide_h), row_lab, FONT_PT, PP_ALIGN.LEFT)
        for j, tr in enumerate(trait_order):
            x = LEFT_W + j * (CELL_W + GAP_X); labs = cell_labels.get((row_lab, tr), [])[:2]
            if len(labs) == 1:
                add_textbox(slide, px_to_ppt(x, canvas_w, slide_w), px_to_ppt(y + CELL_H - LABEL_Y_OFFSET, canvas_h, slide_h),
                            px_to_ppt(CELL_W, canvas_w, slide_w), px_to_ppt(20, canvas_h, slide_h), labs[0], FONT_PT)
            elif len(labs) >= 2:
                add_textbox(slide, px_to_ppt(x, canvas_w, slide_w), px_to_ppt(y + CELL_H - LABEL_Y_OFFSET, canvas_h, slide_h),
                            px_to_ppt(CELL_W/2, canvas_w, slide_w), px_to_ppt(20, canvas_h, slide_h), labs[0], FONT_PT)
                add_textbox(slide, px_to_ppt(x + CELL_W/2, canvas_w, slide_w), px_to_ppt(y + CELL_H - LABEL_Y_OFFSET, canvas_h, slide_h),
                            px_to_ppt(CELL_W/2, canvas_w, slide_w), px_to_ppt(20, canvas_h, slide_h), labs[1], FONT_PT)

    add_textbox(slide, px_to_ppt(x0-8, canvas_w, slide_w), px_to_ppt(y0-36, canvas_h, slide_h),
                px_to_ppt(90, canvas_w, slide_w), px_to_ppt(20, canvas_h, slide_h), "Z score", FONT_PT, PP_ALIGN.LEFT)
    for txt, yy in [(str(ZLIM), y0-5), ("0", y0 + bar_h/2 - 10), (f"-{ZLIM}", y0 + bar_h - 18)]:
        add_textbox(slide, px_to_ppt(x0 + bar_w + 12, canvas_w, slide_w), px_to_ppt(yy, canvas_h, slide_h),
                    px_to_ppt(40, canvas_w, slide_w), px_to_ppt(20, canvas_h, slide_h), txt, FONT_PT, PP_ALIGN.LEFT)
    ppt.save(tmp / "Fig3C.pptx")


def main():
    print("RUNNING Fig3C_LAYOUT_FINAL_20260526")
    print(f"layout: CELL_W={CELL_W}, CELL_H={CELL_H}, GAP_X={GAP_X}, GAP_Y={GAP_Y}, LEFT_W={LEFT_W}, RIGHT_W={RIGHT_W}, VIEW_GAP={VIEW_GAP}, BRAIN_BOX=({BRAIN_BOX_W},{BRAIN_BOX_H}), PPT_W_IN={PPT_W_IN}")
    x = read_phe(); jhu = fetch_jhu()
    cell_files, cell_labels, mapping_rows = {}, {}, []
    for row_lab, cls, meas in row_order:
        xx = x[(x["mri_class"] == cls) & (x["measure"] == meas)]
        if not len(xx): continue
        for tr in trait_order:
            d = xx[xx["trait"] == tr]
            f = cell_dir / f"{clean(row_lab)}_{clean(tr)}.png"
            if cls == "cortex": img, m = draw_cortex(d, f)
            elif cls == "subcortex": img, m = draw_subcortex(d, f)
            else: img, m = draw_white_matter(d, jhu, f)
            cell_files[(row_lab, tr)] = img; cell_labels[(row_lab, tr)] = best_labels(m, 2)
            mapping_rows.append({"row": row_lab, "trait": tr, "n_input_sig": len(d), "n_matched": 0 if m is None else len(m)})
    png, canvas_info, colorbar_info = stitch_no_text(cell_files)
    make_ppt(png, cell_labels, canvas_info, colorbar_info)
    pd.DataFrame(mapping_rows).to_csv(tmp / "Fig3C.mapping_check.csv", index=False)
    shutil.copy2(tmp / "Fig3C.png", out / "Fig3C.png")
    shutil.copy2(tmp / "Fig3C.pptx", out / "Fig3C.pptx")
    shutil.copy2(tmp / "Fig3C.mapping_check.csv", out / "Fig3C.mapping_check.csv")
    layout_report = (
        "RUNNING Fig3C_LAYOUT_FINAL_20260526\n"
        f"CELL_W={CELL_W}\n"
        f"CELL_H={CELL_H}\n"
        f"GAP_X={GAP_X}\n"
        f"GAP_Y={GAP_Y}\n"
        f"LEFT_W={LEFT_W}\n"
        f"RIGHT_W={RIGHT_W}\n"
        f"VIEW_GAP={VIEW_GAP}\n"
        f"BRAIN_BOX_W={BRAIN_BOX_W}\n"
        f"BRAIN_BOX_H={BRAIN_BOX_H}\n"
        f"PPT_W_IN={PPT_W_IN}\n"
    )
    (out / "Fig3C.layout.txt").write_text(layout_report, encoding="utf-8")
    shutil.rmtree(tmp)
    print(out / "Fig3C.png")
    print(out / "Fig3C.pptx")
    print(out / "Fig3C.mapping_check.csv")


if __name__ == "__main__":
    main()
